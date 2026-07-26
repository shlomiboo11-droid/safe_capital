#!/usr/bin/env python3
"""
Build Safe Capital Investor Pitch Deck (.pptx)
Mirrors website/pitch.html — 10 slides + hidden 'internal' slide at end.
RTL Hebrew. Brand colors: Navy #022445, Crimson #984349, Cream #FBF9F6.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand tokens ───────────────────────────────────────────────────────
NAVY        = RGBColor(0x02, 0x24, 0x45)
NAVY_SOFT   = RGBColor(0x1E, 0x3A, 0x5C)
CREAM       = RGBColor(0xFB, 0xF9, 0xF6)
CREAM_SOFT  = RGBColor(0xF5, 0xF3, 0xF0)
CRIMSON     = RGBColor(0x98, 0x43, 0x49)
CRIMSON_SOFT= RGBColor(0xB3, 0x51, 0x58)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x1B, 0x1C, 0x1A)
TEXT_MUTED  = RGBColor(0x43, 0x47, 0x4E)

# ── Layout: 16:9, 13.333" x 7.5" ───────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HEBREW_FONT = 'Heebo'
LATIN_FONT  = 'Inter'

# ──────────────────────────────────────────────────────────────────────
def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs, bg=CREAM):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    # Push to back
    spTree = bg_shape._element.getparent()
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)
    return s

def add_text(slide, x, y, w, h, text, *,
             font_size=18, bold=False, color=TEXT_DARK,
             font=HEBREW_FONT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
             rtl=True, line_height=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_height
        if rtl:
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            pPr.set('rtl', '1')
        run = p.add_run()
        run.text = ln
        f = run.font
        f.name = font
        f.size = Pt(font_size)
        f.bold = bold
        f.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, *, fill=WHITE, outline=False, radius=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not outline:
        shape.line.fill.background()
    return shape

def add_circle(slide, x, y, d, fill=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def slide_header(slide, num, title_text, sub_text=None, *, on_dark=False):
    eyebrow_color = CRIMSON_SOFT if on_dark else CRIMSON
    title_color   = CREAM if on_dark else NAVY
    sub_color     = RGBColor(0xCC, 0xC8, 0xC0) if on_dark else TEXT_MUTED

    # Eyebrow (number)
    add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
             num, font_size=12, bold=True, color=eyebrow_color,
             font=LATIN_FONT)
    # Title
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(1.1),
             title_text, font_size=40, bold=True, color=title_color)
    # Sub
    if sub_text:
        add_text(slide, Inches(0.6), Inches(1.95), Inches(12), Inches(0.7),
                 sub_text, font_size=18, color=sub_color, line_height=1.5)

def page_number(slide, n, total, on_dark=False):
    color = RGBColor(0xBB, 0xB8, 0xB0) if on_dark else TEXT_MUTED
    add_text(slide, Inches(0.6), Inches(7.05), Inches(2), Inches(0.3),
             f'{n} / {total}', font_size=10, color=color, font=LATIN_FONT,
             align=PP_ALIGN.LEFT, rtl=False)
    add_text(slide, Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.3),
             'Safe Capital · סייף קפיטל', font_size=10, color=color,
             font=LATIN_FONT, align=PP_ALIGN.RIGHT)

TOTAL = 10

# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def slide_01_hero(prs):
    s = blank_slide(prs, bg=NAVY)
    # Gradient feel via offset crimson glow (subtle rect)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             'סייף קפיטל · SAFE CAPITAL',
             font_size=13, bold=True, color=CRIMSON_SOFT, font=LATIN_FONT)

    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(3.0),
             'אתם ראשונים ברווח.\nאנחנו מנהלים.\nאתם לא נוגעים בנייר.',
             font_size=64, bold=True, color=CREAM, line_height=1.15)

    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.0),
             'השקעות נדל״ן Flip בבירמינגהם, אלבמה — קצרות, שקופות, בליווי ישראלי מלא.',
             font_size=18, color=RGBColor(0xCC, 0xC8, 0xC0), line_height=1.5)
    page_number(s, 1, TOTAL, on_dark=True)

def slide_02_story(prs):
    s = blank_slide(prs)
    slide_header(s, '02', 'הסיפור שלנו',
                 'ארבעה שותפים, חברים ועמיתים, שבנו ביחד את מה שחסר להם בעצמם בתור משקיעים.')

    # Narrative (right side — wider)
    narrative = (
        'איתן חיפש שלוש שנים את השוק הנכון. הוא נסע לארה״ב, בדק עיר אחר עיר, '
        'ספר אחר ספר, מנהל קבלנים אחר מנהל קבלנים. בסוף עצר בבירמינגהם, אלבמה — '
        'שילוב נדיר של מרווח רווח גבוה, ביקוש יציב, ועוגנים מוסדיים '
        '(UAB Hospital, האוניברסיטה).\n\n'
        'כשחזר ארצה, חיבר את שלומי (CEO), עדי והוטר — שניהם בנויים ליחסי משקיעים. '
        'ביחד החלטנו להציע את אותו מנגנון שאנחנו עצמנו רצינו: עסקה קצרה, '
        'שקיפות מלאה, והמשקיע ראשון בתור על הרווח.'
    )
    add_text(s, Inches(5.5), Inches(2.9), Inches(7.3), Inches(4),
             narrative, font_size=15, color=TEXT_MUTED, line_height=1.6)

    # Team cards (left side, 2x2)
    team = [('שלומי דוד', 'CEO · מנכ״ל', 'ש'),
            ('איתן',       'רכישות · שטח בארה״ב', 'א'),
            ('עדי',        'יחסי משקיעים', 'ע'),
            ('הוטר',       'יחסי משקיעים', 'ה')]
    card_w, card_h = Inches(2.4), Inches(0.95)
    for i, (name, role, initial) in enumerate(team):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 2.55)
        y = Inches(2.95 + row * 1.1)
        add_rect(s, x, y, card_w, card_h, fill=WHITE)
        # Avatar circle
        add_circle(s, x + Inches(1.85), y + Inches(0.175), Inches(0.6), fill=NAVY)
        add_text(s, x + Inches(1.85), y + Inches(0.27), Inches(0.6), Inches(0.4),
                 initial, font_size=18, bold=True, color=CREAM,
                 font=HEBREW_FONT, align=PP_ALIGN.CENTER, rtl=False)
        # Name + role (right-aligned within card, room left for avatar)
        add_text(s, x + Inches(0.1), y + Inches(0.15), Inches(1.65), Inches(0.35),
                 name, font_size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.1), y + Inches(0.5), Inches(1.65), Inches(0.35),
                 role, font_size=10, color=TEXT_MUTED)
    page_number(s, 2, TOTAL)

def slide_03_birmingham(prs):
    s = blank_slide(prs)
    slide_header(s, '03', 'למה דווקא בירמינגהם, אלבמה',
                 'הבחירה לא מקרית. שילוב של מרווח, יציבות ועוגנים מוסדיים שבישראל פשוט לא קיים.')

    # Compare cards
    card_w, card_h = Inches(6), Inches(2.5)
    # Right card: Israel (muted)
    add_rect(s, Inches(6.7), Inches(2.9), card_w, card_h, fill=CREAM_SOFT)
    add_text(s, Inches(6.9), Inches(3.05), Inches(5.6), Inches(0.4),
             'נדל״ן בישראל', font_size=11, bold=True, color=CRIMSON, font=LATIN_FONT)
    add_text(s, Inches(6.9), Inches(3.4), Inches(5.6), Inches(0.7),
             '3–5%   ', font_size=42, bold=True, color=TEXT_MUTED, font=LATIN_FONT,
             rtl=False)
    add_text(s, Inches(6.9), Inches(4.0), Inches(5.6), Inches(0.4),
             'תשואה שנתית', font_size=12, color=TEXT_MUTED)
    bullets_il = '• מחירי כניסה ₪2M+\n• אופק 7–15 שנה\n• חשיפה גיאופוליטית מתמשכת\n• רגולציה כבדה ושוכרים מורכבים'
    add_text(s, Inches(6.9), Inches(4.5), Inches(5.6), Inches(1.5),
             bullets_il, font_size=12, color=TEXT_MUTED, line_height=1.5)

    # Left card: Birmingham (navy gradient feel)
    add_rect(s, Inches(0.6), Inches(2.9), card_w, card_h, fill=NAVY)
    add_text(s, Inches(0.8), Inches(3.05), Inches(5.6), Inches(0.4),
             'Flip בבירמינגהם', font_size=11, bold=True,
             color=CRIMSON_SOFT, font=LATIN_FONT)
    add_text(s, Inches(0.8), Inches(3.4), Inches(5.6), Inches(0.7),
             '18–22%   ', font_size=42, bold=True, color=CREAM,
             font=LATIN_FONT, rtl=False)
    add_text(s, Inches(0.8), Inches(4.0), Inches(5.6), Inches(0.4),
             'לעסקה (6–12 חודשים)', font_size=12,
             color=RGBColor(0xCC, 0xC8, 0xC0))
    bullets_bh = '• כניסה מ-$50K\n• אופק 6–12 חודשים\n• UAB Hospital · 28K עובדים\n• אבטלה 2.2% · מס נכס $919/שנה'
    add_text(s, Inches(0.8), Inches(4.5), Inches(5.6), Inches(1.5),
             bullets_bh, font_size=12, color=RGBColor(0xCC, 0xC8, 0xC0),
             line_height=1.5)

    # Three "How" cards below
    titles = ['Hard Money Lending', 'אופק קצר = הגנה', 'שכונות נבחרות']
    descs = ['מימון מקומי קצר־טווח שמאפשר לרכוש מהר, לשפץ ולמכור — בלי לשעבד הון יקר.',
             '6–12 חודשים. שינויי שוק או גיאופוליטיקה לא מספיקים לפגוע בעסקה אחת.',
             'Over the Mountain — שכונות יוקרה ליד בתי חולים ואוניברסיטאות. ביקוש קשיח.']
    for i in range(3):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(5.7), Inches(4), Inches(1.1), fill=WHITE)
        add_text(s, x + Inches(0.2), Inches(5.8), Inches(3.7), Inches(0.4),
                 titles[i], font_size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), Inches(6.15), Inches(3.7), Inches(0.6),
                 descs[i], font_size=10, color=TEXT_MUTED, line_height=1.4)
    page_number(s, 3, TOTAL)

def slide_04_offer(prs):
    s = blank_slide(prs)
    slide_header(s, '04', 'ההצעה — מה אתה מקבל',
                 'בלי שכבות, בלי כוכביות. שלושה מספרים, מנגנון אחד.')

    # 3 stat cards
    stats = [('$50K',  'השקעה מינימלית', WHITE, NAVY,    TEXT_MUTED),
             ('20%',   'Preferred Return', CRIMSON, CREAM, RGBColor(0xEE, 0xE0, 0xE2)),
             ('6–12',  'חודשים · משך עסקה', WHITE, NAVY,    TEXT_MUTED)]
    for i, (value, label, bg, fg, lbl) in enumerate(stats):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(2.9), Inches(4), Inches(1.7), fill=bg)
        add_text(s, x + Inches(0.3), Inches(3.0), Inches(3.5), Inches(0.9),
                 value, font_size=52, bold=True, color=fg, font=LATIN_FONT,
                 align=PP_ALIGN.RIGHT, rtl=False)
        add_text(s, x + Inches(0.3), Inches(4.05), Inches(3.5), Inches(0.4),
                 label, font_size=13, color=lbl)

    # Waterfall card
    add_rect(s, Inches(0.6), Inches(4.85), Inches(12.15), Inches(2.1), fill=WHITE)
    add_text(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.5),
             'איך מתחלק הרווח', font_size=22, bold=True, color=NAVY)

    steps = [
        ('1', 'המשקיע ראשון בתור — 20% על הקרן שלו',
         'התשואה של 20% מחושבת על הקרן של המשקיע, לא על רווח העסקה.'),
        ('2', 'כל מה שמעל ל-20% למשקיע — אנחנו מקבלים',
         'העסקה צריכה להרוויח מספיק כדי לכסות את ה-20% למשקיע. כל שקל מעבר לזה — שלנו.'),
        ('3', 'אם אין מספיק — המשקיע לוקח הכל, אנחנו אפס',
         'אנחנו לא רואים אגורה לפני שהמשקיע מקבל את ה-20% שלו. Skin in the Game אמיתי.')
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(5.4 + i * 0.5)
        # Step number circle (right side — RTL)
        add_circle(s, Inches(12.0), y - Inches(0.05), Inches(0.4), fill=CREAM_SOFT)
        add_text(s, Inches(12.0), y, Inches(0.4), Inches(0.3),
                 num, font_size=14, bold=True, color=CRIMSON, font=LATIN_FONT,
                 align=PP_ALIGN.CENTER, rtl=False)
        add_text(s, Inches(0.85), y, Inches(11.0), Inches(0.25),
                 title, font_size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.85), y + Inches(0.22), Inches(11.0), Inches(0.3),
                 desc, font_size=10, color=TEXT_MUTED, line_height=1.4)
    page_number(s, 4, TOTAL)

def slide_05_relationship(prs):
    s = blank_slide(prs)
    slide_header(s, '05', 'זה לא One-Shot. זו מערכת יחסים.',
                 'הכוח האמיתי הוא הצבירה. הקרן נשארת אצלנו, מתגלגלת לעסקה הבאה — ומפיקה 20% שוב ושוב.')

    # Cycle stages
    stages = [('1', 'משקיע מעביר $50K', False),
              ('2', 'עסקה: 6–12 חודשים', False),
              ('3', '+20% רווח · קרן נשארת', False),
              ('4', 'העסקה הבאה — אוטומטית', True)]
    cw = Inches(2.8)
    for i, (num, label, loop) in enumerate(stages):
        x = Inches(0.6 + i * 3.15)
        bg = NAVY if loop else WHITE
        fg = CREAM if loop else NAVY
        lbl = RGBColor(0xCC, 0xC8, 0xC0) if loop else TEXT_MUTED
        accent = CRIMSON_SOFT if loop else CRIMSON
        add_rect(s, x, Inches(2.9), cw, Inches(1.3), fill=bg)
        add_text(s, x, Inches(3.0), cw, Inches(0.5),
                 num, font_size=30, bold=True, color=accent, font=LATIN_FONT,
                 align=PP_ALIGN.CENTER, rtl=False)
        add_text(s, x, Inches(3.65), cw, Inches(0.5),
                 label, font_size=12, color=lbl, align=PP_ALIGN.CENTER)
        # arrow between stages
        if i < 3:
            add_text(s, x + cw + Inches(0.03), Inches(3.35), Inches(0.3), Inches(0.4),
                     '◂', font_size=20, color=NAVY, font=LATIN_FONT,
                     align=PP_ALIGN.CENTER, rtl=False)

    # Example card
    add_rect(s, Inches(0.6), Inches(4.55), Inches(12.15), Inches(1.5), fill=CREAM_SOFT)
    add_text(s, Inches(0.85), Inches(4.65), Inches(11.6), Inches(0.4),
             'לדוגמה — ביצוע 18 חודשים', font_size=11, bold=True, color=CRIMSON, font=LATIN_FONT)
    legs = [('השקעת', '$50,000', False),
            ('אחרי עסקה 1', '+$10,000', False),
            ('אחרי עסקה 2', '+$10,000', False),
            ('סך רווח', '$20,000', True)]
    for i, (caption, value, accent) in enumerate(legs):
        x = Inches(0.85 + i * 2.93)
        bg = NAVY if accent else WHITE
        fg_cap = RGBColor(0xCC, 0xC8, 0xC0) if accent else TEXT_MUTED
        fg_val = CREAM if accent else NAVY
        add_rect(s, x, Inches(5.05), Inches(2.75), Inches(0.85), fill=bg)
        add_text(s, x + Inches(0.15), Inches(5.13), Inches(2.5), Inches(0.25),
                 caption, font_size=9, color=fg_cap)
        add_text(s, x + Inches(0.15), Inches(5.4), Inches(2.5), Inches(0.4),
                 value, font_size=18, bold=True, color=fg_val, font=LATIN_FONT,
                 rtl=False, align=PP_ALIGN.RIGHT)

    # 3 options
    add_rect(s, Inches(0.6), Inches(6.15), Inches(12.15), Inches(0.85), fill=WHITE)
    add_text(s, Inches(0.85), Inches(6.22), Inches(11.6), Inches(0.3),
             'בכל סוף עסקה — 3 אפשרויות', font_size=11, bold=True, color=CRIMSON, font=LATIN_FONT)
    opts = 'השאר הכל בחשבון להשקעה הבאה   ·   משוך רווח, השאר קרן   ·   משוך הכל ל-IBAN בישראל'
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.4),
             opts, font_size=12, color=TEXT_MUTED)
    page_number(s, 5, TOTAL)

def slide_06_track(prs):
    s = blank_slide(prs)
    slide_header(s, '06', 'מה כבר עשינו', 'לא מילים — מספרים ונכסים אמיתיים.')

    stats = [('—', 'משקיעים פעילים'),
             ('—', 'הון מנוהל'),
             ('—', 'ROI ממוצע'),
             ('—', 'עסקאות שהושלמו')]
    cw = Inches(2.95)
    for i, (val, label) in enumerate(stats):
        x = Inches(0.6 + i * 3.1)
        add_rect(s, x, Inches(3.0), cw, Inches(1.6), fill=WHITE)
        add_text(s, x + Inches(0.2), Inches(3.2), Inches(2.6), Inches(0.7),
                 val, font_size=44, bold=True, color=NAVY, font=LATIN_FONT, rtl=False)
        add_text(s, x + Inches(0.2), Inches(4.05), Inches(2.6), Inches(0.4),
                 label, font_size=13, color=TEXT_MUTED)

    add_rect(s, Inches(0.6), Inches(5.0), Inches(12.15), Inches(1.3), fill=CREAM_SOFT)
    add_text(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(0.5),
             'המספרים יוצגו ידנית בפגישה.', font_size=18, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)
    page_number(s, 6, TOTAL)

def slide_07_deal(prs):
    s = blank_slide(prs)
    slide_header(s, '07', 'העסקה הקונקרטית — היום',
                 'הנכס שאנחנו מציעים למשקיע בעסקה הנוכחית. הנתונים מהאתר.')

    # Single big card placeholder for the live deal
    add_rect(s, Inches(0.6), Inches(2.9), Inches(12.15), Inches(4), fill=WHITE)

    # Hero strip (placeholder)
    add_rect(s, Inches(0.85), Inches(3.1), Inches(11.6), Inches(1.4), fill=CREAM_SOFT)
    add_text(s, Inches(0.95), Inches(3.2), Inches(11.4), Inches(0.4),
             'גיוס בעיצומו · 12 חודשים · ROI 20%', font_size=11, bold=True,
             color=CRIMSON, font=LATIN_FONT)
    add_text(s, Inches(0.95), Inches(3.55), Inches(11.4), Inches(0.6),
             '1513 Oxmoor Rd', font_size=28, bold=True, color=NAVY,
             font=LATIN_FONT, rtl=False, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.95), Inches(4.05), Inches(11.4), Inches(0.4),
             '1513 Oxmoor Road, Birmingham, AL 35209', font_size=12,
             color=TEXT_MUTED, font=LATIN_FONT, rtl=False, align=PP_ALIGN.RIGHT)

    # 4 metrics
    metrics = [('$475,000', 'מחיר רכישה',     NAVY),
               ('$1,650,000','ARV / מכירה צפויה', NAVY),
               ('$211,461', 'צפי רווח',        CRIMSON),
               ('20%',      'תשואה למשקיע',    CRIMSON)]
    mw = Inches(2.8)
    for i, (val, label, c) in enumerate(metrics):
        x = Inches(0.85 + i * 2.93)
        add_text(s, x, Inches(4.85), mw, Inches(0.4),
                 label, font_size=11, color=TEXT_MUTED)
        add_text(s, x, Inches(5.2), mw, Inches(0.6),
                 val, font_size=24, bold=True, color=c, font=LATIN_FONT, rtl=False)

    add_text(s, Inches(0.85), Inches(6.1), Inches(11.6), Inches(0.6),
             'נכס יוקרתי בשכונת Oxmoor, בירמינגהם — שיפוץ מקיף, מכירה צפויה תוך 12 חודשים, פוטנציאל השבחה משמעותי.',
             font_size=12, color=TEXT_MUTED, line_height=1.5)
    page_number(s, 7, TOTAL)

def slide_08_contract(prs):
    s = blank_slide(prs)
    slide_header(s, '08', 'החוזה — לא מבוסס על אמון בלבד',
                 'מסגרת משפטית ברורה. כל שורה הגיונית, אין שפה מתפלפלת.')

    items = [
        ('סכום, לוח זמנים, תשואה',
         '$50K (או יותר) · עד 12 חודשים · 20% Preferred Return — שחור על גבי לבן.'),
        ('2 חודשי חסד',
         'אם הפרויקט מתעכב — יש לנו חודשיים לסגור. מעבר לכך, המשקיע יכול לכפות buyout בהנחה.'),
        ('המשקיע מוגן אם נכס לא נמכר',
         'אנחנו סופגים הפסד מחיר, לא הקרן. המשקיע יוצא ראשון.'),
        ('LLC על שם המשקיע',
         'לכל עסקה — חברה משפטית נפרדת. שמך על הנייר, יכולת תביעה ישירה.'),
        ('חתימה דיגיטלית או פיזית',
         'כל מה שנוח. עו״ד ישראלי + עו״ד אמריקאי מלווים את התהליך.')
    ]
    # 2 columns × 3 rows (one extra in column)
    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.15)
        y = Inches(2.9 + row * 1.3)
        add_rect(s, x, y, Inches(6), Inches(1.15), fill=WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.4),
                 title, font_size=14, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(5.6), Inches(0.55),
                 desc, font_size=11, color=TEXT_MUTED, line_height=1.5)
    page_number(s, 8, TOTAL)

def slide_09_money_flow(prs):
    s = blank_slide(prs)
    slide_header(s, '09', 'איך הכסף זז', 'פנימה והחוצה. שום דבר מוסתר.')

    # A: inflow
    add_rect(s, Inches(0.6), Inches(2.9), Inches(12.15), Inches(1.3), fill=WHITE)
    add_text(s, Inches(0.85), Inches(3.0), Inches(11.6), Inches(0.35),
             'A · תזרים פנימה', font_size=11, bold=True, color=CRIMSON, font=LATIN_FONT)
    a_text = ('• המשקיע מעביר את הכסף תוך 10 ימים מחתימה  '
              '· חשבון בנק ייעודי לכל פרויקט\n'
              '• המרה ל-USD על המשקיע · עמלות לפי החוזה  '
              '· איחור — לפי הקבוע בחוזה')
    add_text(s, Inches(0.85), Inches(3.35), Inches(11.6), Inches(0.85),
             a_text, font_size=12, color=TEXT_MUTED, line_height=1.5)

    # B: 3 options
    add_text(s, Inches(0.6), Inches(4.4), Inches(12.15), Inches(0.4),
             'B · בסוף עסקה — 3 אפשרויות', font_size=11, bold=True,
             color=CRIMSON, font=LATIN_FONT)
    opts = [('1', 'השאר הכל בחשבון',  'קרן + רווח נשארים, נכנסים ישר לעסקה הבאה. רק מס משולם.', False),
            ('2', 'משוך רווח, השאר קרן', 'הקרן ממשיכה לעבוד. הרווח חוזר ל-IBAN בישראל. מס משולם.', True),
            ('3', 'משוך הכל',          'קרן + רווח חוזרים בחזרה לישראל. מס משולם. בלי שאלות.', False)]
    for i, (num, title, desc, accent) in enumerate(opts):
        x = Inches(0.6 + i * 4.15)
        bg = NAVY if accent else CREAM_SOFT
        fg_num = CRIMSON_SOFT if accent else CRIMSON
        fg_t = CREAM if accent else NAVY
        fg_d = RGBColor(0xCC, 0xC8, 0xC0) if accent else TEXT_MUTED
        add_rect(s, x, Inches(4.8), Inches(4), Inches(1.4), fill=bg)
        add_text(s, x + Inches(0.2), Inches(4.9), Inches(3.6), Inches(0.5),
                 num, font_size=22, bold=True, color=fg_num, font=LATIN_FONT, rtl=False)
        add_text(s, x + Inches(0.2), Inches(5.3), Inches(3.6), Inches(0.35),
                 title, font_size=13, bold=True, color=fg_t)
        add_text(s, x + Inches(0.2), Inches(5.65), Inches(3.6), Inches(0.5),
                 desc, font_size=10, color=fg_d, line_height=1.4)

    # C: management
    add_rect(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.65), fill=WHITE)
    add_text(s, Inches(0.85), Inches(6.4), Inches(11.6), Inches(0.3),
             'C · ניהול', font_size=10, bold=True, color=CRIMSON, font=LATIN_FONT)
    add_text(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.3),
             'עדכון תזרים בכל אבן דרך — דרך פורטל המשקיעים והווצאפ.',
             font_size=11, color=TEXT_MUTED)
    page_number(s, 9, TOTAL)

def slide_10_close(prs):
    s = blank_slide(prs, bg=NAVY)
    add_text(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
             'סגירה · ללא לחץ', font_size=13, bold=True,
             color=CRIMSON_SOFT, font=LATIN_FONT)

    add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(2.8),
             'קח את הזמן\nלחשוב על זה.',
             font_size=72, bold=True, color=CREAM, line_height=1.15)

    add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.4),
             'תישן על זה לילה או שניים. תתייעץ עם מי שצריך.\n'
             'אנחנו לא מחפשים חתימה עכשיו — ניצור איתך קשר בעוד יומיים.',
             font_size=18, color=RGBColor(0xCC, 0xC8, 0xC0), line_height=1.6)

    add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             'Safe Capital · סייף קפיטל', font_size=10,
             color=RGBColor(0xBB, 0xB8, 0xB0), font=LATIN_FONT)

def slide_11_internal(prs):
    """HIDDEN preparation slide — internal only."""
    s = blank_slide(prs, bg=RGBColor(0xF3, 0xE8, 0xE9))
    add_rect(s, 0, 0, SLIDE_W, Inches(0.5), fill=CRIMSON)
    add_text(s, Inches(0.6), Inches(0.08), Inches(12), Inches(0.4),
             '⚠  INTERNAL — DO NOT SHOW  ·  דף הכנה לשלומי',
             font_size=14, bold=True, color=CREAM, font=LATIN_FONT,
             align=PP_ALIGN.CENTER, rtl=False)

    add_text(s, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
             'תשובות מוכנות', font_size=12, bold=True, color=CRIMSON, font=LATIN_FONT)
    add_text(s, Inches(0.6), Inches(1.05), Inches(12), Inches(0.55),
             '7 ההתנגדויות הנפוצות', font_size=30, bold=True, color=NAVY)

    objections = [
        ('1. עיכוב בלוח זמנים — השיפוץ לוקח יותר זמן',
         'חוזית: 2 חודשי חסד. אם עובר 3 חודשים מהדדליין — המשקיע יכול לכפות buyout בהנחה.'),
        ('2. הנכס לא נמכר / נתקע בשוק',
         'נכסים דומים נמכרים תוך 25 ימים מקסימום. אם נתקע — אנחנו מורידים מחיר וסופגים. הקרן מוגנת.'),
        ('3. איך אני יודע שהכסף שלי באמת הולך לשיפוץ?',
         'חשבון בנק ייעודי לכל פרויקט. שקיפות מלאה דרך פורטל המשקיעים — רואים כל תזרים.'),
        ('4. מה אם יש בעיה משפטית או נזק באתר הבנייה?',
         'אנחנו נמנעים מנכסים בעייתיים מראש. ביטוח בנייה מלא מכסה פציעות או תקריות.'),
        ('5. מסים וציות — האם זה בטוח למשקיע ישראלי?',
         'אנחנו מטפלים בכל המסמכים. מחברים לרו״ח מומחה — דיווח מס בשתי המדינות.'),
        ('6. מה אם השוק קורס במהלך הפרויקט?',
         'עסקאות קצרות = תנודות שוק לא מספיקות לפגוע. קונים מתחת לשוק. המשקיע ראשון על הרווח.'),
        ('7. מה אם תיעלמו / משהו יקרה לצוות?',
         'יחסים אישיים (חברים, מילואים יחד) · LLC על שם המשקיע · יכולת לתבוע · עתידית: ערבות בנקאית.')
    ]
    for i, (title, desc) in enumerate(objections):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.15)
        y = Inches(1.85 + row * 1.3)
        add_rect(s, x, y, Inches(6), Inches(1.2), fill=WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.4),
                 title, font_size=12, bold=True, color=CRIMSON)
        add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(5.6), Inches(0.7),
                 desc, font_size=10, color=TEXT_MUTED, line_height=1.4)

    add_text(s, Inches(0.6), Inches(7.15), Inches(12), Inches(0.3),
             'השקופית הזו לא להצגה — דף הכנה בלבד.',
             font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Mark slide as hidden in PowerPoint
    sld = s._element
    show = etree.SubElement(sld, qn('p:show'))
    sld.set('show', '0')

# ══════════════════════════════════════════════════════════════════════
def main():
    prs = make_prs()
    slide_01_hero(prs)
    slide_02_story(prs)
    slide_03_birmingham(prs)
    slide_04_offer(prs)
    slide_05_relationship(prs)
    slide_06_track(prs)
    slide_07_deal(prs)
    slide_08_contract(prs)
    slide_09_money_flow(prs)
    slide_10_close(prs)
    slide_11_internal(prs)

    out = '/Users/shlomidavid/claudecode/safe_capital/.claude/worktrees/eager-bhabha-2fcae1/pitch-deck.pptx'
    prs.save(out)
    print(f'OK → {out}')

if __name__ == '__main__':
    main()

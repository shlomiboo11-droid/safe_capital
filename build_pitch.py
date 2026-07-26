"""
Safe Capital — Pitch Deck Builder
מצגת השקעה לפגישה אישית עם משקיע
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── BRAND COLORS ──────────────────────────────────────────────
NAVY     = RGBColor(0x02, 0x24, 0x45)
NAVY2    = RGBColor(0x1E, 0x3A, 0x5C)
CRIMSON  = RGBColor(0x98, 0x43, 0x49)
CREAM    = RGBColor(0xFB, 0xF8, 0xF3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x43, 0x47, 0x4E)
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xF0)
GOLD     = RGBColor(0xC8, 0xA9, 0x6E)   # accent warm

# ── SLIDE SIZE: 16:9 ─────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # fully blank

# ── HELPERS ───────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_name="Arial", font_size=24, bold=False,
                color=WHITE, align=PP_ALIGN.RIGHT,
                wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # RTL paragraph
    set_rtl(p)
    return txBox

def add_multiline_textbox(slide, lines, x, y, w, h,
                          font_name="Arial", font_size=18,
                          bold=False, color=WHITE,
                          align=PP_ALIGN.RIGHT,
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        set_rtl(p)
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        if isinstance(line, dict):
            run.text = line.get('text', '')
            run.font.name = font_name
            run.font.size = Pt(line.get('size', font_size))
            run.font.bold = line.get('bold', bold)
            run.font.color.rgb = line.get('color', color)
        else:
            run.text = line
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox

def set_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')

def slide_number_tag(slide, num, total=12, bg_dark=True):
    txt_color = RGBColor(0xAA, 0xBB, 0xCC) if bg_dark else RGBColor(0x74, 0x77, 0x7F)
    add_textbox(slide, f"{num} / {total}",
                W - Inches(1.4), H - Inches(0.45), Inches(1.2), Inches(0.35),
                font_name="Arial", font_size=11, color=txt_color,
                align=PP_ALIGN.CENTER)

def divider_line(slide, y, color=NAVY2, width=W - Inches(2), x=Inches(1)):
    bar = add_rect(slide, x, y, width, Pt(1.5), fill_color=color)
    return bar

def stat_card(slide, number, label, x, y, w=Inches(2.4), h=Inches(1.5),
              num_color=CRIMSON, bg=WHITE):
    card = add_rect(slide, x, y, w, h, fill_color=bg)
    # Number
    add_textbox(slide, number, x + Inches(0.1), y + Inches(0.12),
                w - Inches(0.2), Inches(0.7),
                font_name="Arial", font_size=28, bold=True,
                color=num_color, align=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, label, x + Inches(0.1), y + Inches(0.75),
                w - Inches(0.2), Inches(0.6),
                font_name="Arial", font_size=12, bold=False,
                color=GRAY, align=PP_ALIGN.CENTER)

def section_header(slide, title, subtitle=None, bg_dark=True):
    """Big title + optional subtitle in upper area"""
    tc = WHITE if bg_dark else NAVY
    add_textbox(slide, title,
                Inches(0.7), Inches(0.35), W - Inches(1.4), Inches(0.75),
                font_name="Arial", font_size=32, bold=True,
                color=tc, align=PP_ALIGN.RIGHT)
    if subtitle:
        sc = RGBColor(0xAA, 0xBB, 0xCC) if bg_dark else GRAY
        add_textbox(slide, subtitle,
                    Inches(0.7), Inches(1.0), W - Inches(1.4), Inches(0.45),
                    font_name="Arial", font_size=16, bold=False,
                    color=sc, align=PP_ALIGN.RIGHT)

def bg_fill(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

# Big company name
add_textbox(s, "סייף קפיטל",
            Inches(1), Inches(1.6), W - Inches(2), Inches(1.4),
            font_name="Arial", font_size=72, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

# English name smaller
add_textbox(s, "SAFE CAPITAL",
            Inches(1), Inches(2.9), W - Inches(2), Inches(0.5),
            font_name="Arial", font_size=18, bold=False,
            color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)

# Crimson divider
add_rect(s, Inches(4.5), Inches(3.55), Inches(4.3), Pt(3), fill_color=CRIMSON)

# Tagline
add_textbox(s, "נדל\"ן אמריקאי. ניהול ישראלי. אתם ראשונים ברווח.",
            Inches(1), Inches(3.75), W - Inches(2), Inches(0.7),
            font_name="Arial", font_size=20, bold=False,
            color=CREAM, align=PP_ALIGN.CENTER)

# Sub
add_textbox(s, "השקעות Flip בבירמינגהם, אלבמה — עבור משקיעים ישראלים",
            Inches(1.5), Inches(4.5), W - Inches(3), Inches(0.5),
            font_name="Arial", font_size=14, bold=False,
            color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)

# Bottom strip
add_rect(s, Inches(0), H - Inches(0.6), W, Inches(0.6), fill_color=NAVY2)
add_textbox(s, "מצגת לפגישה אישית — חסוי",
            Inches(0.5), H - Inches(0.55), W - Inches(1), Inches(0.45),
            font_name="Arial", font_size=11, color=RGBColor(0x6B, 0x7A, 0x8D),
            align=PP_ALIGN.CENTER)
slide_number_tag(s, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — הסיפור שלנו
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, CREAM)

section_header(s, "הסיפור שלנו", "למה סייף קפיטל?", bg_dark=False)
add_rect(s, Inches(0.7), Inches(1.35), Inches(4), Pt(2), fill_color=CRIMSON)

# Story text (right column)
story_lines = [
    "הכל התחיל עם איתן.",
    "",
    "לפני כמה שנים, איתן החל לחקור את שוק הנדל\"ן האמריקאי.",
    "הוא ביקר בעשרות ערים, קרא מאות דוחות שוק, ובסוף",
    "הגיע לאותה מסקנה: בירמינגהם, אלבמה.",
    "",
    "שם גילה שבשכונות היוקרה — Homewood, Mountain Brook —",
    "כמעט ואין בנייה חדשה. כשיוצא בית ישן למכירה, יש תור",
    "של קונים. הפוטנציאל לרווח מהשבחה היה ברור.",
    "",
    "לאחר עסקה ראשונה מוצלחת, שלומי, עדי וחוטר הצטרפו.",
    "ביחד — פתחנו את סייף קפיטל לשותפים."
]

add_multiline_textbox(s, story_lines,
                      W - Inches(6.5), Inches(1.6), Inches(5.8), Inches(4.5),
                      font_name="Arial", font_size=14, color=GRAY,
                      align=PP_ALIGN.RIGHT, line_spacing=2)

# Team cards (left side)
team = [
    ("שלומי דוד", "מייסד ומנכ\"ל"),
    ("איתן", "מנהל רכישות"),
    ("עדי", "יחסי משקיעים"),
    ("חוטר", "יחסי משקיעים"),
]
card_w = Inches(2.6)
card_h = Inches(1.3)
card_x = Inches(0.4)
for i, (name, role) in enumerate(team):
    cy = Inches(1.7) + i * (card_h + Inches(0.2))
    add_rect(s, card_x, cy, card_w, card_h, fill_color=WHITE)
    # Name
    add_textbox(s, name, card_x + Inches(0.15), cy + Inches(0.18),
                card_w - Inches(0.3), Inches(0.5),
                font_name="Arial", font_size=15, bold=True,
                color=NAVY, align=PP_ALIGN.RIGHT)
    # Role
    add_textbox(s, role, card_x + Inches(0.15), cy + Inches(0.6),
                card_w - Inches(0.3), Inches(0.4),
                font_name="Arial", font_size=12, color=CRIMSON,
                align=PP_ALIGN.RIGHT)

# Bottom tagline
add_rect(s, Inches(0), H - Inches(0.7), W, Inches(0.7), fill_color=NAVY)
add_textbox(s, "אנחנו לא חברה אנונימית — אנחנו ארבעה שותפים שמשקיעים מהכיס שלנו בכל עסקה",
            Inches(0.5), H - Inches(0.65), W - Inches(1), Inches(0.55),
            font_name="Arial", font_size=13, bold=True,
            color=CREAM, align=PP_ALIGN.CENTER)
slide_number_tag(s, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — השיטה שלנו
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

section_header(s, "השיטה שלנו — לא הימור", "6 שלבים לפני כל עסקה", bg_dark=True)
add_rect(s, Inches(0.7), Inches(1.35), Inches(4), Pt(2), fill_color=CRIMSON)

steps = [
    ("1", "איתור נכס ישן מתחת לשוק",       "המתווכים שלנו מאתרים לפני שהנכס יוצא לשוק הפתוח"),
    ("2", "בדיקת היתכנות בנייה",            "סטרוקטורה, היתרים, בעיות נסתרות — הכל נבדק"),
    ("3", "תכנית השבחה + חישוב רווח",       "מחשבים עלויות מלאות ובודקים אם הרווח מצדיק"),
    ("4", "שמאי + השוואת Comps",            "שמאי מעריך ARV. משווים לנכסים שנמכרו לאחרונה"),
    ("5", "אישור גוף מימון חיצוני",         "Hard Money Lender בודק את התוכנית באופן בלתי תלוי"),
    ("6", "יציאה לעבודה — עד 12 חודשים",   "רק אחרי שכל הצ'ק-ליסט מאושר — רוכשים ומתחילים"),
]

col_w = Inches(6.0)
col_x = W - col_w - Inches(0.5)

for i, (num, title, desc) in enumerate(steps):
    row_y = Inches(1.65) + i * Inches(0.87)
    # Number circle
    add_rect(s, col_x - Inches(0.05), row_y, Inches(0.45), Inches(0.45), fill_color=CRIMSON)
    add_textbox(s, num, col_x - Inches(0.05), row_y,
                Inches(0.45), Inches(0.45),
                font_name="Arial", font_size=14, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(s, title,
                col_x + Inches(0.5), row_y,
                col_w - Inches(0.6), Inches(0.35),
                font_name="Arial", font_size=14, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)
    # Desc
    add_textbox(s, desc,
                col_x + Inches(0.5), row_y + Inches(0.35),
                col_w - Inches(0.6), Inches(0.38),
                font_name="Arial", font_size=11, color=RGBColor(0x8A, 0xA4, 0xCC),
                align=PP_ALIGN.RIGHT)

# Left column — professionals
pro_x = Inches(0.4)
pro_w = Inches(5.5)
add_textbox(s, "אנשי המקצוע שלנו",
            pro_x, Inches(1.65), pro_w, Inches(0.4),
            font_name="Arial", font_size=16, bold=True,
            color=CREAM, align=PP_ALIGN.RIGHT)

pros = [
    "מתווכים מקומיים — מביאים עסקאות לפני השוק",
    "קבלני שיפוץ קבועים — אותם אנשים, אותה איכות",
    "שמאים מוסמכים — ARV אובייקטיבי",
    "עורך דין אמריקאי + עורך דין ישראלי",
    "רואה חשבון שמכיר את שני הצדדים",
    "Hard Money Lender — אישור בלתי תלוי",
]
for i, pro in enumerate(pros):
    add_textbox(s, f"• {pro}",
                pro_x, Inches(2.25) + i * Inches(0.73),
                pro_w, Inches(0.55),
                font_name="Arial", font_size=13,
                color=CREAM, align=PP_ALIGN.RIGHT)

slide_number_tag(s, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — למה בירמינגהם
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, CREAM)

section_header(s, "למה בירמינגהם, אלבמה?", "שוק שבחרנו בכוונה — לא במקרה", bg_dark=False)
add_rect(s, Inches(0.7), Inches(1.35), Inches(5), Pt(2), fill_color=CRIMSON)

# Stat cards row
stats = [
    ("879,000", "תושבי מטרו"),
    ("+5,700", "תושבים חדשים בשנה"),
    ("2.2%", "שיעור אבטלה"),
    ("$919", "מס נכס שנתי\n(vs. $2,400 ארצי)"),
    ("$12.1B", "השפעה כלכלית\nשנתית של UAB"),
]
card_w = Inches(2.35)
gap = Inches(0.18)
start_x = Inches(0.4)
for i, (num, lbl) in enumerate(stats):
    cx = start_x + i * (card_w + gap)
    stat_card(s, num, lbl, cx, Inches(1.65), w=card_w, h=Inches(1.4))

# Insight box
add_rect(s, Inches(0.4), Inches(3.25), W - Inches(0.8), Inches(1.35), fill_color=NAVY)
add_textbox(s, "🏘  ההזדמנות: בשכונות Homewood, Mountain Brook, Vestavia Hills — כמעט אין בנייה חדשה",
            Inches(0.6), Inches(3.32), W - Inches(1.2), Inches(0.5),
            font_name="Arial", font_size=15, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(s, "כשיוצא בית ישן למכירה — זה נדיר. יש תור קונים. אנחנו קונים, משפצים, מוכרים לאותו תור.",
            Inches(0.6), Inches(3.78), W - Inches(1.2), Inches(0.7),
            font_name="Arial", font_size=13,
            color=CREAM, align=PP_ALIGN.RIGHT)

# Neighborhoods
hoods = [
    ("Homewood", "בתי ספר מצוינים, בוטיקים, ביקוש קבוע"),
    ("Mountain Brook", "היוקרתית ביותר — מחירים גבוהים במיוחד"),
    ("Vestavia Hills", "משפחות, ביקוש גובר, מיקום אסטרטגי"),
    ("Hoover", "צמיחה מהירה, ביקוש לבתים מאובזרים"),
]
hood_w = Inches(2.9)
for i, (name, desc) in enumerate(hoods):
    hx = Inches(0.4) + i * (hood_w + Inches(0.22))
    add_rect(s, hx, Inches(4.75), hood_w, Inches(1.05), fill_color=WHITE)
    add_textbox(s, name, hx + Inches(0.1), Inches(4.82),
                hood_w - Inches(0.2), Inches(0.38),
                font_name="Arial", font_size=13, bold=True,
                color=NAVY, align=PP_ALIGN.RIGHT)
    add_textbox(s, desc, hx + Inches(0.1), Inches(5.15),
                hood_w - Inches(0.2), Inches(0.55),
                font_name="Arial", font_size=11,
                color=GRAY, align=PP_ALIGN.RIGHT)

# UAB note
add_textbox(s, "UAB — המעסיק הגדול באלבמה: 28,000 עובדים, המדורג #1 בבתי חולים 12 שנה ברציפות",
            Inches(0.5), H - Inches(0.55), W - Inches(1), Inches(0.45),
            font_name="Arial", font_size=11,
            color=GRAY, align=PP_ALIGN.CENTER)
slide_number_tag(s, 4, bg_dark=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — ההצעה
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

section_header(s, "ההצעה — ברורה ושקופה", None, bg_dark=True)
add_rect(s, Inches(0.7), Inches(1.1), Inches(3.5), Pt(2), fill_color=CRIMSON)

# 3 big stat cards
main_stats = [
    ("$50,000", "השקעה מינימלית"),
    ("20%", "Preferred Return"),
    ("6–12 חודשים", "משך עסקה"),
]
card_w = Inches(3.7)
gap = Inches(0.4)
total_w = 3 * card_w + 2 * gap
sx = (float(W) - total_w) / 2
for i, (num, lbl) in enumerate(main_stats):
    cx = sx + i * (card_w + gap)
    add_rect(s, cx, Inches(1.4), card_w, Inches(1.5), fill_color=NAVY2)
    add_textbox(s, num, cx, Inches(1.5), card_w, Inches(0.75),
                font_name="Arial", font_size=36, bold=True,
                color=CRIMSON, align=PP_ALIGN.CENTER)
    add_textbox(s, lbl, cx, Inches(2.2), card_w, Inches(0.5),
                font_name="Arial", font_size=14,
                color=CREAM, align=PP_ALIGN.CENTER)

# Preferred Return explanation — 2 scenarios
add_textbox(s, "איך זה עובד:",
            Inches(0.6), Inches(3.1), W - Inches(1.2), Inches(0.38),
            font_name="Arial", font_size=16, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)

# Scenario A
add_rect(s, Inches(0.5), Inches(3.55), Inches(5.8), Inches(1.5), fill_color=NAVY2)
add_textbox(s, "עסקה הצליחה טוב (26% רווח)",
            Inches(0.65), Inches(3.62), Inches(5.5), Inches(0.38),
            font_name="Arial", font_size=14, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(s, "אתם מקבלים: 20%    |    אנחנו מקבלים: 6%",
            Inches(0.65), Inches(3.98), Inches(5.5), Inches(0.38),
            font_name="Arial", font_size=13,
            color=CREAM, align=PP_ALIGN.RIGHT)
add_textbox(s, "✅ תרחיש A",
            Inches(0.65), Inches(4.35), Inches(2), Inches(0.35),
            font_name="Arial", font_size=11,
            color=RGBColor(0x6B, 0xC9, 0x8A), align=PP_ALIGN.RIGHT)

# Scenario B
add_rect(s, Inches(6.9), Inches(3.55), Inches(5.8), Inches(1.5), fill_color=NAVY2)
add_textbox(s, "עסקה הצליחה סבירות (17% רווח)",
            Inches(7.05), Inches(3.62), Inches(5.5), Inches(0.38),
            font_name="Arial", font_size=14, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(s, "אתם מקבלים: 17% (הכל)    |    אנחנו: 0",
            Inches(7.05), Inches(3.98), Inches(5.5), Inches(0.38),
            font_name="Arial", font_size=13,
            color=CREAM, align=PP_ALIGN.RIGHT)
add_textbox(s, "✅ תרחיש B",
            Inches(7.05), Inches(4.35), Inches(2), Inches(0.35),
            font_name="Arial", font_size=11,
            color=RGBColor(0x6B, 0xC9, 0x8A), align=PP_ALIGN.RIGHT)

# Key message
add_rect(s, Inches(0.5), Inches(5.25), W - Inches(1), Inches(0.9), fill_color=CRIMSON)
add_textbox(s, "האינטרסים שלנו מיושרים לחלוטין — אנחנו לא נרוויח אם אתם לא תרוויחו קודם",
            Inches(0.6), Inches(5.35), W - Inches(1.2), Inches(0.65),
            font_name="Arial", font_size=16, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

slide_number_tag(s, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — מודל היחסים
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, CREAM)

section_header(s, "לא עסקה — מערכת יחסים", "הכסף שלך עובד שוב ושוב", bg_dark=False)
add_rect(s, Inches(0.7), Inches(1.35), Inches(4), Pt(2), fill_color=CRIMSON)

# Cycle diagram (text-based)
cycle_items = [
    (Inches(5.0), Inches(1.8),  "השקעה ראשונה\n$50,000"),
    (Inches(9.5), Inches(2.8),  "עסקה מסתיימת\n(6-12 חודש)"),
    (Inches(8.5), Inches(4.8),  "רווח $10,000\nעובר אליכם"),
    (Inches(2.5), Inches(4.8),  "הקרן נשארת\nבחשבון"),
    (Inches(1.5), Inches(2.8),  "עסקה הבאה\nמיידית"),
]
for cx, cy, txt in cycle_items:
    add_rect(s, cx, cy, Inches(2.8), Inches(0.85), fill_color=NAVY)
    add_textbox(s, txt, cx + Inches(0.1), cy + Inches(0.08),
                Inches(2.6), Inches(0.75),
                font_name="Arial", font_size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

# Example
add_rect(s, Inches(0.4), Inches(5.15), W - Inches(0.8), Inches(0.85), fill_color=NAVY2)
add_textbox(s, "דוגמה: $50,000 × שתי עסקאות × 20% = $20,000 רווח בשנה | לעומת ~$1,500 בפיקדון ישראלי",
            Inches(0.6), Inches(5.25), W - Inches(1.2), Inches(0.65),
            font_name="Arial", font_size=14, bold=True,
            color=CREAM, align=PP_ALIGN.CENTER)

# 3 options
opts = [
    ("🔄  השאר הכל", "קרן + רווח ממשיכים לעסקה הבאה"),
    ("💸  קח את הרווח", "הרווח עובר לישראל, הקרן ממשיכה"),
    ("🏦  צא לגמרי", "קרן + רווח חוזרים לחשבון שלך"),
]
opt_w = Inches(3.8)
for i, (title, desc) in enumerate(opts):
    ox = Inches(0.5) + i * (opt_w + Inches(0.25))
    add_rect(s, ox, Inches(6.1), opt_w, Inches(1.0), fill_color=WHITE)
    add_textbox(s, title, ox + Inches(0.1), Inches(6.17),
                opt_w - Inches(0.2), Inches(0.38),
                font_name="Arial", font_size=13, bold=True,
                color=NAVY, align=PP_ALIGN.RIGHT)
    add_textbox(s, desc, ox + Inches(0.1), Inches(6.5),
                opt_w - Inches(0.2), Inches(0.45),
                font_name="Arial", font_size=11,
                color=GRAY, align=PP_ALIGN.RIGHT)

slide_number_tag(s, 6, bg_dark=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Track Record
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

section_header(s, "מה שכבר עשינו", "הוכחה — לא רק מילים", bg_dark=True)
add_rect(s, Inches(0.7), Inches(1.35), Inches(3.5), Pt(2), fill_color=CRIMSON)

# Top stats
tr_stats = [
    ("[X]",    "עסקאות מושלמות"),
    ("[X]",    "משקיעים פעילים"),
    ("[$X]",   "הון מנוהל מצטבר"),
    ("[X%]",   "ROI ממוצע מאומת"),
]
ts_w = Inches(2.9)
ts_gap = Inches(0.35)
tsx = (float(W) - (4 * ts_w + 3 * ts_gap)) / 2
for i, (num, lbl) in enumerate(tr_stats):
    cx = tsx + i * (ts_w + ts_gap)
    add_rect(s, cx, Inches(1.6), ts_w, Inches(1.3), fill_color=NAVY2)
    add_textbox(s, num, cx, Inches(1.7), ts_w, Inches(0.6),
                font_name="Arial", font_size=30, bold=True,
                color=CRIMSON, align=PP_ALIGN.CENTER)
    add_textbox(s, lbl, cx, Inches(2.25), ts_w, Inches(0.45),
                font_name="Arial", font_size=12,
                color=CREAM, align=PP_ALIGN.CENTER)

# Property cards
deals = [
    ("206 Mountain Ave", "Homewood, AL", "הושלם בהצלחה"),
    ("3425 Ridge Dell Cir", "Birmingham, AL", "הושלם בהצלחה"),
    ("1513 Oxmoor Rd", "Birmingham, AL", "בעבודה"),
]
pc_w = Inches(3.8)
pc_gap = Inches(0.38)
pc_x0 = (float(W) - (3 * pc_w + 2 * pc_gap)) / 2
for i, (addr, city, status) in enumerate(deals):
    px = pc_x0 + i * (pc_w + pc_gap)
    add_rect(s, px, Inches(3.15), pc_w, Inches(2.8), fill_color=NAVY2)
    # Status badge
    badge_color = CRIMSON if "בעבודה" in status else RGBColor(0x2A, 0x7A, 0x4A)
    add_rect(s, px + Inches(0.15), Inches(3.25), Inches(1.5), Inches(0.32), fill_color=badge_color)
    add_textbox(s, status,
                px + Inches(0.15), Inches(3.25), Inches(1.5), Inches(0.32),
                font_name="Arial", font_size=10, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    # Address
    add_textbox(s, addr, px + Inches(0.15), Inches(3.65),
                pc_w - Inches(0.3), Inches(0.45),
                font_name="Arial", font_size=14, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)
    add_textbox(s, city, px + Inches(0.15), Inches(4.08),
                pc_w - Inches(0.3), Inches(0.35),
                font_name="Arial", font_size=12,
                color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.RIGHT)
    # Placeholder for data
    add_textbox(s, "[מחיר רכישה / מכירה / ROI]",
                px + Inches(0.15), Inches(4.48),
                pc_w - Inches(0.3), Inches(0.4),
                font_name="Arial", font_size=11,
                color=RGBColor(0xCC, 0xCC, 0x44), align=PP_ALIGN.RIGHT)

add_textbox(s, "* נתוני Track Record לאימות סופי לפני הפגישה",
            Inches(0.5), H - Inches(0.5), W - Inches(1), Inches(0.4),
            font_name="Arial", font_size=10,
            color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)
slide_number_tag(s, 7)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — הנכס הספציפי (Oxmoore)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, CREAM)

section_header(s, "העסקה שלפניכם", "1513 Oxmoor Rd, Birmingham AL 35209  |  Homewood", bg_dark=False)
add_rect(s, Inches(0.7), Inches(1.35), Inches(5), Pt(2), fill_color=CRIMSON)

# Property transformation
add_rect(s, Inches(0.4), Inches(1.6), Inches(3.5), Inches(1.45), fill_color=NAVY)
add_textbox(s, "לפני שיפוץ",
            Inches(0.55), Inches(1.68), Inches(3.2), Inches(0.35),
            font_name="Arial", font_size=12, bold=True,
            color=CREAM, align=PP_ALIGN.RIGHT)
add_textbox(s, "2 חד' | 2 מקלחות | 1,600 sqft",
            Inches(0.55), Inches(2.0), Inches(3.2), Inches(0.35),
            font_name="Arial", font_size=13,
            color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(s, "←",
            Inches(3.95), Inches(1.85), Inches(0.5), Inches(0.5),
            font_name="Arial", font_size=24, bold=True,
            color=CRIMSON, align=PP_ALIGN.CENTER)
add_rect(s, Inches(4.55), Inches(1.6), Inches(3.5), Inches(1.45), fill_color=NAVY2)
add_textbox(s, "אחרי שיפוץ",
            Inches(4.7), Inches(1.68), Inches(3.2), Inches(0.35),
            font_name="Arial", font_size=12, bold=True,
            color=CREAM, align=PP_ALIGN.RIGHT)
add_textbox(s, "5 חד' | 4 מקלחות | 3,000 sqft",
            Inches(4.7), Inches(2.0), Inches(3.2), Inches(0.35),
            font_name="Arial", font_size=13, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)

# Pro Forma (left side)
pf_lines = [
    {"text": "מחיר רכישה:          $475,000", "size": 13, "color": GRAY},
    {"text": "עלות שיפוץ (קבלן):   $505,000", "size": 13, "color": GRAY},
    {"text": "עלויות נלוות:           $22,600", "size": 13, "color": GRAY},
    {"text": "עלויות החזקה (8 חו'):    $7,848", "size": 13, "color": GRAY},
    {"text": "מימון (HML):           $71,402", "size": 13, "color": GRAY},
    {"text": "─────────────────────────────", "size": 11, "color": RGBColor(0xCC, 0xCC, 0xCC)},
    {"text": "סך עלות כוללת:    $1,082,000", "size": 14, "bold": True, "color": NAVY},
    {"text": " ", "size": 8, "color": GRAY},
    {"text": "מחיר מכירה (ARV):  $1,350,000", "size": 13, "color": GRAY},
    {"text": "הוצאות מכירה:        ($81,000)", "size": 13, "color": GRAY},
    {"text": "─────────────────────────────", "size": 11, "color": RGBColor(0xCC, 0xCC, 0xCC)},
    {"text": "רווח נקי משוער:      $181,452", "size": 15, "bold": True, "color": CRIMSON},
]
add_multiline_textbox(s, pf_lines,
                      Inches(0.4), Inches(3.2), Inches(5.8), Inches(3.8),
                      font_name="Arial", font_size=13,
                      color=GRAY, align=PP_ALIGN.RIGHT)

# Right side — waterfall + comps
# Investor return box
add_rect(s, Inches(6.6), Inches(3.2), Inches(6.1), Inches(1.55), fill_color=NAVY)
add_textbox(s, "השקעת $50,000?",
            Inches(6.8), Inches(3.28), Inches(5.7), Inches(0.38),
            font_name="Arial", font_size=14, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(s, "קיבלת בחזרה $60,000 תוך 8 חודשים",
            Inches(6.8), Inches(3.65), Inches(5.7), Inches(0.38),
            font_name="Arial", font_size=13,
            color=CREAM, align=PP_ALIGN.RIGHT)
add_textbox(s, "רווח: $10,000 = 20%",
            Inches(6.8), Inches(4.02), Inches(5.7), Inches(0.5),
            font_name="Arial", font_size=20, bold=True,
            color=CRIMSON, align=PP_ALIGN.RIGHT)

# Comps
add_textbox(s, "Comps — נכסים דומים שנמכרו באותו אזור:",
            Inches(6.6), Inches(4.95), Inches(6.1), Inches(0.38),
            font_name="Arial", font_size=13, bold=True,
            color=NAVY, align=PP_ALIGN.RIGHT)
comps = [
    ("133 E Edgewood Dr", "$1,350,000", "5/4 | 3,027 sqft"),
    ("1511 Grove Pl",     "$1,357,000", "5/3 | 2,824 sqft"),
    ("1410 Ardsley Pl",   "$1,530,000", "4/4 | 2,998 sqft"),
]
for i, (addr, price, specs) in enumerate(comps):
    cy = Inches(5.4) + i * Inches(0.58)
    add_rect(s, Inches(6.6), cy, Inches(6.1), Inches(0.5), fill_color=LIGHT_BG)
    add_textbox(s, addr, Inches(6.7), cy + Inches(0.07),
                Inches(3), Inches(0.35),
                font_name="Arial", font_size=12,
                color=NAVY, align=PP_ALIGN.RIGHT)
    add_textbox(s, price, Inches(9.8), cy + Inches(0.07),
                Inches(1.5), Inches(0.35),
                font_name="Arial", font_size=12, bold=True,
                color=CRIMSON, align=PP_ALIGN.CENTER)
    add_textbox(s, specs, Inches(11.3), cy + Inches(0.07),
                Inches(1.3), Inches(0.35),
                font_name="Arial", font_size=10,
                color=GRAY, align=PP_ALIGN.LEFT)

add_textbox(s, "✅ ה-ARV שלנו ($1,350,000) הוא השמרני ביותר — יש מרווח ביטחון",
            Inches(6.6), Inches(7.1), Inches(6.1), Inches(0.35),
            font_name="Arial", font_size=11, bold=True,
            color=RGBColor(0x2A, 0x7A, 0x4A), align=PP_ALIGN.RIGHT)

slide_number_tag(s, 8, bg_dark=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — החוזה
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

section_header(s, "המסגרת המשפטית", "אמון טוב יותר כשיש מסמך מאחוריו", bg_dark=True)
add_rect(s, Inches(0.7), Inches(1.35), Inches(4), Pt(2), fill_color=CRIMSON)

contract_items = [
    ("✅", "סכום ההשקעה וההתחייבות", "מוגדר מראש. סכום קבוע לעסקה ספציפית."),
    ("✅", "לוח זמנים מחייב", "תאריך יעד לסיום הפרויקט — כתוב בחוזה."),
    ("✅", "Preferred Return — 20%", "המשקיע מקבל 20% על הקרן לפני שהמייסדים מקבלים דבר."),
    ("✅", "חסד — 2 חודשים", "אם הפרויקט מתעכב — יש 2 חודשי חסד מעבר ל-deadline."),
    ("✅", "Buyout בהנחה",  "אם העיכוב חורג מ-3 חודשים — המשקיע יכול לכפות יציאה בהנחה."),
    ("✅", "המשקיע מוגן אם הנכס לא נמכר", "אנחנו סופגים הפסדים מהפחתת מחיר. הקרן של המשקיע בראש."),
    ("✅", "מבנה LLC", "הנכס רשום על שם ה-LLC. שמך מופיע בניירות הרשמיים."),
    ("✅", "אופציות יציאה מוגדרות", "ברור מתי ואיך ניתן לצאת מהעסקה."),
]

col_w = Inches(5.8)
col_h = len(contract_items) * Inches(0.66)
left_x = W - col_w - Inches(0.5)

for i, (icon, title, desc) in enumerate(contract_items):
    iy = Inches(1.6) + i * Inches(0.66)
    # Icon
    add_textbox(s, icon, left_x - Inches(0.05), iy,
                Inches(0.4), Inches(0.38),
                font_name="Arial", font_size=14,
                color=RGBColor(0x6B, 0xC9, 0x8A), align=PP_ALIGN.CENTER)
    # Title
    add_textbox(s, title, left_x + Inches(0.45), iy,
                col_w - Inches(0.5), Inches(0.3),
                font_name="Arial", font_size=13, bold=True,
                color=WHITE, align=PP_ALIGN.RIGHT)
    # Desc
    add_textbox(s, desc, left_x + Inches(0.45), iy + Inches(0.3),
                col_w - Inches(0.5), Inches(0.3),
                font_name="Arial", font_size=11,
                color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.RIGHT)

# Right column visual
add_rect(s, Inches(0.4), Inches(1.6), Inches(5.5), Inches(5.3), fill_color=NAVY2)
add_textbox(s, "📄",
            Inches(1.5), Inches(2.3), Inches(3), Inches(1.2),
            font_name="Arial", font_size=60,
            color=CRIMSON, align=PP_ALIGN.CENTER)
add_textbox(s, "חוזה שותפות",
            Inches(0.6), Inches(3.6), Inches(5.2), Inches(0.5),
            font_name="Arial", font_size=20, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "LLC ייעודי לכל עסקה",
            Inches(0.6), Inches(4.1), Inches(5.2), Inches(0.4),
            font_name="Arial", font_size=14,
            color=CREAM, align=PP_ALIGN.CENTER)
add_textbox(s, "שמך כמשקיע מופיע\nבניירות הרשמיים",
            Inches(0.6), Inches(4.65), Inches(5.2), Inches(0.7),
            font_name="Arial", font_size=13,
            color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)

slide_number_tag(s, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — לוגיסטיקה פיננסית
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, CREAM)

section_header(s, "איך הכסף זז", "שקיפות מלאה — כל שקל מתועד", bg_dark=False)
add_rect(s, Inches(0.7), Inches(1.35), Inches(4.5), Pt(2), fill_color=CRIMSON)

# Inflow
add_rect(s, Inches(0.4), Inches(1.6), Inches(5.8), Inches(2.1), fill_color=NAVY)
add_textbox(s, "תזרים פנימה — מהמשקיע אלינו",
            Inches(0.55), Inches(1.68), Inches(5.5), Inches(0.38),
            font_name="Arial", font_size=14, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)
inflow_lines = [
    "• תוך 10 ימים מחתימת החוזה",
    "• חשבון בנק ייעודי לכל פרויקט בנפרד",
    "• קנס במקרה של איחור [לפי חוזה]",
]
for i, line in enumerate(inflow_lines):
    add_textbox(s, line,
                Inches(0.55), Inches(2.12) + i * Inches(0.42),
                Inches(5.5), Inches(0.38),
                font_name="Arial", font_size=12,
                color=CREAM, align=PP_ALIGN.RIGHT)

# 3 options cards
opts_data = [
    ("🔄  השאר הכל",
     "קרן + רווח ממשיכים לעסקה הבאה",
     "מקסימום צבר"),
    ("💸  קח את הרווח",
     "הרווח ($10K) לישראל, הקרן ממשיכה",
     "הכנסה שוטפת"),
    ("🏦  צא לגמרי",
     "קרן + רווח חוזרים לחשבון ישראלי",
     "גמישות מלאה"),
]
opt_w = Inches(3.9)
opt_gap = Inches(0.3)
opt_x0 = Inches(0.4)
for i, (title, desc, tag) in enumerate(opts_data):
    ox = opt_x0 + i * (opt_w + opt_gap)
    add_rect(s, ox, Inches(3.9), opt_w, Inches(2.2), fill_color=WHITE)
    add_rect(s, ox, Inches(3.9), opt_w, Inches(0.38), fill_color=NAVY)
    add_textbox(s, tag, ox + Inches(0.1), Inches(3.92),
                opt_w - Inches(0.2), Inches(0.3),
                font_name="Arial", font_size=11, bold=True,
                color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)
    add_textbox(s, title, ox + Inches(0.1), Inches(4.38),
                opt_w - Inches(0.2), Inches(0.45),
                font_name="Arial", font_size=14, bold=True,
                color=NAVY, align=PP_ALIGN.RIGHT)
    add_textbox(s, desc, ox + Inches(0.1), Inches(4.85),
                opt_w - Inches(0.2), Inches(0.7),
                font_name="Arial", font_size=12,
                color=GRAY, align=PP_ALIGN.RIGHT)

# Tax note
add_rect(s, Inches(0.4), Inches(6.25), W - Inches(0.8), Inches(0.6), fill_color=NAVY2)
add_textbox(s, "מס: ישולם פעם בשנה בסיוע רו\"ח המתמחה בשני הצדדים — ישראל + ארה\"ב",
            Inches(0.6), Inches(6.33), W - Inches(1.2), Inches(0.45),
            font_name="Arial", font_size=13,
            color=CREAM, align=PP_ALIGN.CENTER)

slide_number_tag(s, 10, bg_dark=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — סגירה
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, NAVY)

# Crimson top stripe
add_rect(s, Inches(0), Inches(0), W, Inches(0.25), fill_color=CRIMSON)

add_textbox(s, "תקח את הזמן שלך",
            Inches(1), Inches(1.0), W - Inches(2), Inches(1.1),
            font_name="Arial", font_size=56, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(4.5), Inches(2.2), Inches(4.3), Pt(2), fill_color=CRIMSON)

main_text = [
    "אין צורך להחליט עכשיו.",
    "",
    "ישן על זה לילה-שניים.",
    "דבר עם מי שצריך.",
    "תחזור עם שאלות.",
]
add_multiline_textbox(s, main_text,
                      Inches(2), Inches(2.5), W - Inches(4), Inches(2.2),
                      font_name="Arial", font_size=22,
                      color=CREAM, align=PP_ALIGN.CENTER, line_spacing=4)

add_textbox(s, "ניצור איתך קשר בעוד יומיים לשמוע מה בראש שלך.",
            Inches(2), Inches(4.8), W - Inches(4), Inches(0.55),
            font_name="Arial", font_size=16,
            color=RGBColor(0x8A, 0xA4, 0xCC), align=PP_ALIGN.CENTER)

# Contact info
add_rect(s, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.8), fill_color=NAVY2)
add_textbox(s, "📱 054-7828550   |   ✉️ safecapital2024@gmail.com",
            Inches(3.6), Inches(5.6), Inches(6.1), Inches(0.55),
            font_name="Arial", font_size=14,
            color=CREAM, align=PP_ALIGN.CENTER)

add_textbox(s, "אנחנו לא מחפשים \"כן\" מהיר — אנחנו מחפשים שותפים לטווח ארוך.",
            Inches(1), Inches(6.5), W - Inches(2), Inches(0.5),
            font_name="Arial", font_size=14, bold=True,
            color=CRIMSON, align=PP_ALIGN.CENTER)

slide_number_tag(s, 11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7.5 — הכנה פנימית (HIDDEN — slide 12)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg_fill(s, RGBColor(0x1A, 0x1A, 0x2E))  # Very dark — clearly different

# Warning header
add_rect(s, Inches(0), Inches(0), W, Inches(0.7), fill_color=RGBColor(0x8B, 0x00, 0x00))
add_textbox(s, "⚠️  פנימי — לא להציג למשקיע  ⚠️  Cmd+Shift+H",
            Inches(0.3), Inches(0.08), W - Inches(0.6), Inches(0.5),
            font_name="Arial", font_size=16, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "הכנה — תשובות להתנגדויות",
            Inches(0.6), Inches(0.85), W - Inches(1.2), Inches(0.55),
            font_name="Arial", font_size=24, bold=True,
            color=WHITE, align=PP_ALIGN.RIGHT)

objections = [
    ("1. השיפוץ מתעכב",
     "2 חודשי חסד. אחרי 3 חודשים מה-deadline — buyout בהנחה."),
    ("2. הנכס לא נמכר",
     "Comps: נמכרים תוך 25 ימים. אם נתקע — אנחנו מורידים מחיר. הקרן מוגנת."),
    ("3. איך אני יודע שהכסף הולך לשיפוץ?",
     "חשבון בנק ייעודי לפרויקט. פורטל משקיעים — כל תזרים גלוי."),
    ("4. תאונה/בעיה משפטית?",
     "ביטוח בנייה מלא. אנחנו נמנעים מנכסים עם סיכון משפטי מראש."),
    ("5. מסים — לא רוצה להסתבך",
     "מחברים עם רו\"ח ישראלי-אמריקאי. הוא מטפל בכל הדיווח."),
    ("6. השוק קורס?",
     "6-12 חודשים = חשיפה מינימלית. קונים 30% מתחת לשוק. משקיע מקבל ראשון."),
    ("7. מה אם תיעלמו?",
     "LLC — שמך על הנכס. יכולת תביעה. בעתיד: ערבות בנקאית."),
]
for i, (q, a) in enumerate(objections):
    iy = Inches(1.6) + i * Inches(0.77)
    if i < 4:
        ox = Inches(0.4)
        ow = Inches(5.8)
    else:
        ox = Inches(0.4) + (i-4) % 3 * Inches(4.2)
        ow = Inches(4.0)

    ox = Inches(0.4) if i % 2 == 0 else Inches(6.9)
    ow = Inches(5.8)
    iy = Inches(1.6) + (i // 2) * Inches(0.85)

    add_rect(s, ox, iy, ow, Inches(0.75), fill_color=RGBColor(0x2A, 0x2A, 0x45))
    add_textbox(s, q, ox + Inches(0.1), iy + Inches(0.05),
                ow - Inches(0.2), Inches(0.3),
                font_name="Arial", font_size=12, bold=True,
                color=RGBColor(0xFF, 0xCC, 0x44), align=PP_ALIGN.RIGHT)
    add_textbox(s, a, ox + Inches(0.1), iy + Inches(0.35),
                ow - Inches(0.2), Inches(0.35),
                font_name="Arial", font_size=11,
                color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.RIGHT)

add_textbox(s, "כלל זהב: שלב את התשובות באופן טבעי בשקופיות 8-9 — אל תציג שקופית זו",
            Inches(0.4), H - Inches(0.5), W - Inches(0.8), Inches(0.4),
            font_name="Arial", font_size=11,
            color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ── SAVE ────────────────────────────────────────────────────────
output_path = "/Users/shlomidavid/claudecode/safe_capital/safe_capital_pitch.pptx"
prs.save(output_path)
print(f"✅ נשמר: {output_path}")
print(f"   שקופיות: {len(prs.slides)}")
